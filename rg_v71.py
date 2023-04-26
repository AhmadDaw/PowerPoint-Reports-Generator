from pptx import *
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CATEGORY_TYPE
from pptx.enum.text import PP_ALIGN
import numpy as np
import pandas as pd
from tkinter import *
from extract_cell_name import extract_cel_nam
from tkinter import filedialog
from datetime import datetime
# -------------------------------------------------------

root=Tk()
root.title("Reports Generator")

title2=Label(root, text="Reports Generator")

eng=Label(root, text="Engineer:")
e_eng=Entry(root, width=20)

stitle=Label(root, text="Site:")
e_stitle=Entry(root, width=20)

load_2_title=Label(root, text="Load 2G file:")
load_3_title=Label(root, text="Load 3G file:")
load_4_title=Label(root, text="Load 4G file:")

current_state2=Label(root, text="-")
status = Label(root, text = "Version: 071 - Coded by: Ahmad Dawara", bd=2, relief=SUNKEN, anchor = E)
# --------------------------------------------------------
fp2='x'
fp3='x'
fp4='x'
def opn_2g():
    global fp2
    fp2=filedialog.askopenfilename()
def opn_3g():
    global fp3
    fp3=filedialog.askopenfilename()
def opn_4g():
    global fp4
    fp4=filedialog.askopenfilename()

b2=Button(root, text="Browse", command=opn_2g)
b3=Button(root, text="Browse", command=opn_3g)
b4=Button(root, text="Browse", command=opn_4g)
# -------------------------------------------------------
CheckVar1 = IntVar()
c1 = Checkbutton(root, text = "Don't Use Default KPIs", variable = CheckVar1, onvalue = 1, offvalue = 0)
CheckVar2 = IntVar()
c2 = Checkbutton(root, text = "Don't Clean Cell Name", variable = CheckVar2, onvalue = 1, offvalue = 0)
# -------------------------------------------------------
def proc():
    dt=datetime.today().strftime('%Y-%m-%d')
    prs = Presentation()
    #prs.slide_width=Inches(16)
    #prs.slide_height=Inches(9)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle = slide.placeholders[1]
    title.text = "Integration Report for site: "+e_stitle.get()
    subtitle.text = "Prepared by: "+e_eng.get()+"\n"+'Date: '+str(dt)
    # -------------------------------------------------------
    if fp2!='x':
        slide_2g_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_2g_layout)
        title = slide.shapes.title
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text = "2G"

        df= pd.read_excel(fp2)
        if (CheckVar1.get() == 0):
            df=extract_cel_nam(df)
            xaxis_col_name=df.columns[0]
            cell_column_name=df.columns[3]
            cells_list=df[cell_column_name].unique()
            df.rename(columns={'cTCH Traffic (Erl)': 'TCH Traffic (Erl)',
                            'SD_AVAIL (%)':'SDCCH Availability (%)',
                            'aCall Setup Success (%)':'Call Setup Success Rate (%)',
                             'Radio_Handover Success Rate (%)':'Handover Success Rate (%)' }
                            , inplace=True)
                        
            kpis_list=['SDCCH Availability (%)','TCH Traffic (Erl)','Call Setup Success Rate (%)','Handover Success Rate (%)']
        
        
        if (CheckVar1.get() == 1):
            df_config= pd.read_excel('rg-config.xlsx')
            cell_col=int(df_config.iloc[0,0])
            print(cell_col)
            dates_col=int(df_config.iloc[0,1])
            cell_col=cell_col-1
            dates_col=dates_col-1
            kpis_list=list(df_config['2G'].unique())
            kpis_list.pop() 
            print(kpis_list)
            df.iloc[:,2:] = df.iloc[:,2:].replace(np.nan, 0)
            xaxis_col_name=df.columns[dates_col]
            cell_column_name=df.columns[cell_col]
            cells_list=df[cell_column_name].unique()

        if (CheckVar2.get() == 0):
            df.iloc[:,4:] = df.iloc[:,4:].replace('NIL', 0)
        
        cells_list.sort()
        i=0
        print(df)

        o=0
        k_len=len(kpis_list)
        number_of_cells=len(cells_list)
        f=0
        for k in kpis_list:
            for cell in cells_list:
                
                dfc=df[df[cell_column_name]==cell]
                dfc = dfc.sort_values(xaxis_col_name)
                dfc[xaxis_col_name] = dfc[xaxis_col_name].map(str)
                #dfc.to_csv('dfc.csv')
                xaxis_lst=list(dfc[xaxis_col_name])
                #print(list(dfc[str(k)]))
                yaxis = tuple(list(dfc[str(k)]))
                if i==0:
                    s2_r=prs.slide_layouts[5]
                    s2=prs.slides.add_slide(s2_r)
                    s2_title=s2.shapes.title
                    s2_title.text=kpis_list[o]
                    s2_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    title_para = s2.shapes.title.text_frame.paragraphs[0]
                    title_para.font.size = Pt(26)

                chart_data = ChartData()
                chart_data.categories = xaxis_lst
                chart_data.add_series(cell, yaxis)

                if (number_of_cells==1):
                    if i==0:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)

                elif (number_of_cells==2):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    else:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)

                elif (number_of_cells==3):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)
                elif (number_of_cells==4 or number_of_cells==8):
                    if i==0:
                        x, y, cx, cy = Inches(3.2), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)                

                elif (number_of_cells==5 or number_of_cells==10):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(6.4), Inches(1.2), Inches(3.2), Inches(3.2) 
                    elif i==4:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2) 

                elif (number_of_cells==6 or number_of_cells==12):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(3.2), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(6.4), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2) 
                    elif i==4:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2) 
                    elif i==5:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2) 

                chart = s2.shapes.add_chart(
                    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
                ).chart

                chart.has_legend = False
                chart.series[0].smooth = True
                chart.has_title = True
                chart.chart_title.text_frame.text = str(cell)
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(15)

                i=i+1
                if (number_of_cells==1) and (i==1):
                    i=0
                elif (number_of_cells==2) and (i==2):
                    i=0
                elif (number_of_cells==3) and (i==3):
                    i=0
                elif (number_of_cells==4 and i==4) or (number_of_cells==8 and i==4):
                    i=0
                elif (number_of_cells==5 and i==5) or (number_of_cells==10 and i==5):
                    i=0
                elif (number_of_cells==6 and i==6) or (number_of_cells==12 and i==6):
                    i=0

                category_axis = chart.category_axis
                category_axis.category_type == XL_CATEGORY_TYPE.TIME_SCALE
                category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
                category_axis.tick_labels.font.size = Pt(12)

                value_axis = chart.value_axis
                mx=max(yaxis)+5
                value_axis.maximum_scale = mx
                value_axis.minimum_scale = 0
                value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
                
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(12)


            o=o+1
        current_state2.config(text='2G Done.')
    # --------------------------------------------------------------
    if fp3!='x':
        slide_3g_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_3g_layout)
        title = slide.shapes.title
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text = "3G"

        df= pd.read_excel(fp3)
        if (CheckVar1.get() == 0):
            df=extract_cel_nam(df)
            xaxis_col_name=df.columns[0]
            cell_column_name=df.columns[3]
            cells_list=df[cell_column_name].unique()
            df.rename(columns={'VS.RAB.AMR.Erlang.cell (Erl)': 'CS Traffic (Erl)',
                        'Radio_CS Inter-Rat Handover Success Rate (%)': 'CS Inter-Rat Handover SR (%)',
                        '{Upgrade}RRC Setup Success Ratio (Service) (%)': 'RRC Setup Success Ratio (Service) (%)',
                        '{Upgrade}RRC Setup Success Ratio (Other) (%)':'RRC Setup Success Ratio (Other) (%)',
                        'VS.HSDPA.UE.Mean.Cell (None)':'Mean HSDPA UE Cell',
                        'Radio_CS RAB Assignment Success Rate (%)':'CS RAB Assignment Success Rate (%)',
                        'Radio_PS RAB Assignment Success Rate (%)':'PS RAB Assignment Success Rate (%)',
                        'HSDPA MAC-d MegaByte (MB)':'HSDPA Data Volume (MByte)'}
                        , inplace=True)
            kpis_list=['CS Traffic (Erl)','RRC Setup Success Ratio (Service) (%)','CS RAB Assignment Success Rate (%)','PS RAB Assignment Success Rate (%)','HSDPA RLC Throughput (kbit/s)','HSDPA Data Volume (MByte)','Mean HSDPA UE Cell']
        
        
        if (CheckVar1.get() == 1):
            df_config= pd.read_excel('rg-config.xlsx')
            cell_col=int(df_config.iloc[0,0])
            dates_col=int(df_config.iloc[0,1])
            cell_col=cell_col-1
            dates_col=dates_col-1            
            kpis_list=list(df_config['3G'].unique())
            kpis_list.pop() 
            print(kpis_list)            
            df.iloc[:,2:] = df.iloc[:,2:].replace(np.nan, 0)
            xaxis_col_name=df.columns[dates_col]
            cell_column_name=df.columns[cell_col]
            cells_list=df[cell_column_name].unique()

        if (CheckVar2.get() == 0):
            df.iloc[:,4:] = df.iloc[:,4:].replace('NIL', 0)
                
        cells_list.sort()
        number_of_cells=len(cells_list)
        #number_of_cells=number_of_cells/2
        f=0
        i=0
        print(df)

        o=0
        fx=0
        k_len=len(kpis_list)
        for k in kpis_list:
            for cell in cells_list:
                dfc=df[df[cell_column_name]==cell]
                dfc = dfc.sort_values(xaxis_col_name)
                dfc[xaxis_col_name] = dfc[xaxis_col_name].map(str)
                xaxis_lst=list(dfc[xaxis_col_name])
                yaxis = tuple(list(dfc[str(k)]))
                if i==0:
                    s2_r=prs.slide_layouts[5]
                    s2=prs.slides.add_slide(s2_r)
                    s2_title=s2.shapes.title
                    s2_title.text=kpis_list[o]
                    s2_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    title_para = s2.shapes.title.text_frame.paragraphs[0]
                    title_para.font.size = Pt(26)


                chart_data = ChartData()
                chart_data.categories = xaxis_lst
                chart_data.add_series(cell, yaxis)

                if (number_of_cells==1):
                    if i==0:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)

                elif (number_of_cells==2):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    else:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)

                elif (number_of_cells==3):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)
                elif (number_of_cells==4 or number_of_cells==8):
                    if i==0:
                        x, y, cx, cy = Inches(3.2), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)                

                elif (number_of_cells==5 or number_of_cells==10):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(6.4), Inches(1.2), Inches(3.2), Inches(3.2) 
                    elif i==4:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2) 

                elif (number_of_cells==6 or number_of_cells==12):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(3.2), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(6.4), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2) 
                    elif i==4:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2) 
                    elif i==5:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2) 


                chart = s2.shapes.add_chart(
                    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
                ).chart

                chart.has_legend = False
                #chart.legend.include_in_layout = False
                chart.series[0].smooth = True
                chart.has_title = True
                chart.chart_title.text_frame.text = str(cell)
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(15)

                i=i+1
                if (number_of_cells==1) and (i==1):
                    i=0
                elif (number_of_cells==2) and (i==2):
                    i=0
                elif (number_of_cells==3) and (i==3):
                    i=0
                elif (number_of_cells==4 and i==4) or (number_of_cells==8 and i==4):
                    i=0
                elif (number_of_cells==5 and i==5) or (number_of_cells==10 and i==5):
                    i=0
                elif (number_of_cells==6 and i==6) or (number_of_cells==12 and i==6):
                    i=0

                category_axis = chart.category_axis
                category_axis.category_type == XL_CATEGORY_TYPE.TIME_SCALE
                category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
                #category_axis.tick_labels.font.italic = True
                category_axis.tick_labels.font.size = Pt(12)

                value_axis = chart.value_axis
                mx=max(yaxis)+5
                value_axis.maximum_scale = mx
                value_axis.minimum_scale = 0
                value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
                
                tick_labels = value_axis.tick_labels
                #tick_labels.font.bold = True
                tick_labels.font.size = Pt(12)
            o=o+1
        current_state2.config(text='3G Done.')
    # --------------------------------------------------------------
    if fp4!='x':
        slide_4g_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_4g_layout)
        title = slide.shapes.title
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text = "4G"

        df= pd.read_excel(fp4)
        if (CheckVar1.get() == 0):
            df=extract_cel_nam(df)
            xaxis_col_name=df.columns[0]
            cell_column_name=df.columns[3]
            cells_list=df[cell_column_name].unique()
            if 'L.Thrp.bits.DL (bit)' in df.columns:
                df['Data Volume (GByte)'] = df['L.Thrp.bits.DL (bit)']/(1024*1024*1024*8)
            if 'L.Thrp.bits.DL.LastTTI (bit)' in df.columns:
                df['User Throughput (Mbps)'] = ((df['L.Thrp.bits.DL (bit)']-df['L.Thrp.bits.DL.LastTTI (bit)'])/df['L.Thrp.Time.DL.RmvLastTTI (ms)']*1000)/1000/1000
            df.iloc[:,4:] = df.iloc[:,4:].replace(np.nan, 0)

            df.rename(columns={'L.ChMeas.PRB.DL.Used.Avg (None)': 'Avg Used PRB DL',
                                'L.Traffic.User.Max (None)': 'Max User Number',
                                'RCC_SetupSuccessRate (Signaling) (%)':'RRC Setup Success Rate (Signaling) (%)',
                                'E-RAB_Setup_Success_Rate (All) (%)':'E-RAB Setup Success Rate (All) (%)'}
                            , inplace=True)
            kpis_list=['Data Volume (GByte)','Avg Used PRB DL','Max User Number','User Throughput (Mbps)','RRC Setup Success Rate (Service) (%)','RRC Setup Success Rate (Signaling) (%)','E-RAB Setup Success Rate (All) (%)']
            
        if (CheckVar1.get() == 1):
            df_config= pd.read_excel('rg-config.xlsx')
            cell_col=int(df_config.iloc[0,0])
            dates_col=int(df_config.iloc[0,1])
            cell_col=cell_col-1
            dates_col=dates_col-1            
            kpis_list=list(df_config['4G'].unique())
            kpis_list.pop() 
            print(kpis_list)            
            df.iloc[:,2:] = df.iloc[:,2:].replace(np.nan, 0)
            xaxis_col_name=df.columns[dates_col]
            cell_column_name=df.columns[cell_col]
            cells_list=df[cell_column_name].unique()

        if (CheckVar2.get() == 0):
            df.iloc[:,4:] = df.iloc[:,4:].replace('NIL', 0)

        cells_list.sort()
        number_of_cells=len(cells_list)
        f=0
        
        i=0
        print(df)

        o=0
        k_len=len(kpis_list)
        for k in kpis_list:
            for cell in cells_list:
                dfc=df[df[cell_column_name]==cell]
                dfc = dfc.sort_values(xaxis_col_name)
                dfc[xaxis_col_name] = dfc[xaxis_col_name].map(str)
                xaxis_lst=list(dfc[xaxis_col_name])
                yaxis = tuple(list(dfc[str(k)]))
                if i==0:
                    s2_r=prs.slide_layouts[5]
                    s2=prs.slides.add_slide(s2_r)
                    s2_title=s2.shapes.title
                    s2_title.text=kpis_list[o]
                    s2_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    title_para = s2.shapes.title.text_frame.paragraphs[0]
                    title_para.font.size = Pt(26)
                
                chart_data = ChartData()
                chart_data.categories = xaxis_lst
                chart_data.add_series(cell, yaxis)

                if (number_of_cells==1):
                    if i==0:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)

                elif (number_of_cells==2):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    else:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)

                elif (number_of_cells==3):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)
                elif (number_of_cells==4):
                    if i==0:
                        x, y, cx, cy = Inches(3.2), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)                

                elif (number_of_cells==5):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(6.4), Inches(1.2), Inches(3.2), Inches(3.2) 
                    elif i==4:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2) 

                elif (number_of_cells==6 or number_of_cells==12):
                    if i==0:
                        x, y, cx, cy = Inches(0), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==1:  
                        x, y, cx, cy = Inches(3.2), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==2:
                        x, y, cx, cy = Inches(6.4), Inches(1.2), Inches(3.2), Inches(3.2)
                    elif i==3:
                        x, y, cx, cy = Inches(0), Inches(4.3), Inches(3.2), Inches(3.2) 
                    elif i==4:
                        x, y, cx, cy = Inches(3.2), Inches(4.3), Inches(3.2), Inches(3.2) 
                    elif i==5:
                        x, y, cx, cy = Inches(6.4), Inches(4.3), Inches(3.2), Inches(3.2)  

                chart = s2.shapes.add_chart(
                    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
                ).chart

                chart.has_legend = False
                #chart.legend.include_in_layout = False
                chart.series[0].smooth = True
                chart.has_title = True
                chart.chart_title.text_frame.text = str(cell)
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(15)

                i=i+1
                if (number_of_cells==1) and (i==1):
                    i=0
                elif (number_of_cells==2) and (i==2):
                    i=0
                elif (number_of_cells==3) and (i==3):
                    i=0
                elif (number_of_cells==4 and i==4) or (number_of_cells==8 and i==4):
                    i=0
                elif (number_of_cells==5 and i==5) or (number_of_cells==10 and i==5):
                    i=0
                elif (number_of_cells==6 and i==6) or (number_of_cells==12 and i==6):
                    i=0

                category_axis = chart.category_axis
                category_axis.category_type == XL_CATEGORY_TYPE.TIME_SCALE
                category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
                #category_axis.tick_labels.font.italic = True
                category_axis.tick_labels.font.size = Pt(12)

                value_axis = chart.value_axis
                mx=max(yaxis)+5
                value_axis.maximum_scale = mx
                value_axis.minimum_scale = 0
                value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
                
                tick_labels = value_axis.tick_labels
                #tick_labels.font.bold = True
                tick_labels.font.size = Pt(12)
            o=o+1
        current_state2.config(text='4G Done.')
    # --------------------------------------------------------------
    
    prs.save(e_stitle.get()+'.pptx')
    current_state2.config(text='Done.')

# -----------------------------------------------------------------
b=Button(root, text="Generate Report", command=proc)

title2.grid(row = 0, column = 1, pady=5, padx=5)
stitle.grid(row = 1, column = 0, pady=5, padx=5)
e_stitle.grid(row = 1, column = 1, pady=5, padx=5)

eng.grid(row = 2, column = 0, pady=5, padx=5)
e_eng.grid(row = 2, column = 1, pady=5, padx=5)

c1.grid(row = 3, column = 1, pady=5, padx=10) 
c2.grid(row = 4, column = 1, pady=5, padx=10) 

load_2_title.grid(row = 5, column = 0, pady=5, padx=5)
load_3_title.grid(row = 6, column = 0, pady=5, padx=5)
load_4_title.grid(row = 7, column = 0, pady=5, padx=5)

b2.grid(row = 5, column = 1, pady=5) 
b3.grid(row = 6, column = 1, pady=5) 
b4.grid(row = 7, column = 1, pady=5) 

current_state2.grid(row = 8, column = 1, pady=5, padx=5)

b.grid(row = 9, column = 1, pady=5)  

status.grid(row=12, column=0, columnspan=6, sticky=W+E)
# ----------------------------------------------------------
root.mainloop()


